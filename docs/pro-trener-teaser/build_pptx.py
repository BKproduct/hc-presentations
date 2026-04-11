#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Генерация PPTX для тизера Health Copilot × Pro Trener (python-pptx)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(13, 27, 62)
TEAL = RGBColor(13, 148, 136)
AMBER = RGBColor(245, 158, 11)
TEXT = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)


def _blank_slide(prs):
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name == "Blank":
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[-1])


def _bg_navy(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def _add_text(slide, left, top, width, height, text, size, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return box


def build(path: Path) -> None:
    prs = Presentation()
    prs.core_properties.title = "Health Copilot — тизер для Pro Trener"
    prs.core_properties.author = "AI-Automation.Studio"

    # 1
    s = _blank_slide(prs)
    _bg_navy(s)
    _add_text(s, Inches(0.55), Inches(1.1), Inches(8.9), Inches(1.1), "Health Copilot", 44, True, TEXT, PP_ALIGN.LEFT)
    _add_text(s, Inches(0.55), Inches(2.15), Inches(8.9), Inches(0.5), "Тизер для команды Pro Trener", 18, False, MUTED)
    _add_text(s, Inches(0.55), Inches(2.85), Inches(8.9), Inches(1.4),
              "Связка тренировок, анализов, генетики, питания и носимых — в одном слое. Без «ещё одного дашборда».", 15, False, TEXT)
    bar = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(4.35), Inches(2.6), Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    _add_text(s, Inches(0.55), Inches(4.6), Inches(8.9), Inches(0.35), "Telegram · мультиагенты · источники · трекеры", 12, False, MUTED)

    slides_copy = [
        ("Утро в Telegram", "Короткий брифинг: что смотреть сегодня, что синхронизировать с тренировкой и питанием."),
        ("Бриф к визиту врачу", "Структурированный конспект: симптомы, динамика лабораторий, вопросы к специалисту, контекст добавок."),
        ("Мультиагентная система", "Отдельные агенты для лаборатории, геномики, носимых, нутрициологии и координации — один ответ пользователю."),
        ("«Мозг» внутри", "Оркестрация: память, политика безопасности, уровни доказательности, ссылки на первоисточники."),
        ("Порядок из хаоса", "Фото анализов и еды из Telegram раскладываются по папкам: анализы, генетика, таблетки, питание, заметки."),
        ("Геном: несколько этапов", "QC → аннотация → приоритизация для обсуждения → связка с фенотипом и образом жизни (не диагноз)."),
        ("Трекеры", "Oura / Apple Health / Whoop — в одном контексте со сном, нагрузкой и лабораторией."),
        ("Источники", "PubMed и клинические руководства — с грейдингом A/B/C. Без «магии нейросети»."),
        ("Что получает клиент", "Меньше трения между залом, нутрициологией и врачом. Яснее приоритеты между сессиями."),
    ]

    for title, body in slides_copy:
        s = _blank_slide(prs)
        _bg_navy(s)
        _add_text(s, Inches(0.55), Inches(0.55), Inches(8.9), Inches(0.55), title, 28, True, TEAL)
        _add_text(s, Inches(0.55), Inches(1.25), Inches(8.9), Inches(3.8), body, 17, False, TEXT)
        _add_text(s, Inches(0.55), Inches(5.15), Inches(8.9), Inches(0.45),
                  "Не медицинское устройство. Не заменяет очную консультацию.", 11, False, MUTED)

    s = _blank_slide(prs)
    _bg_navy(s)
    _add_text(s, Inches(0.55), Inches(1.9), Inches(8.9), Inches(0.9), "Следующий шаг", 26, True, AMBER)
    _add_text(s, Inches(0.55), Inches(2.95), Inches(8.9), Inches(1.2), "Демо и пилот под ваш бренд — 20 минут.", 18, False, TEXT)
    _add_text(s, Inches(0.55), Inches(4.25), Inches(8.9), Inches(0.6), "Telegram: @BKproduct", 22, True, TEAL)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print("Wrote", path)


if __name__ == "__main__":
    build(Path(__file__).resolve().parent / "Health_Copilot_Pro_Trener.pptx")
